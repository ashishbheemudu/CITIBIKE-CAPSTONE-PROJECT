#!/usr/bin/env python3
"""
NYC Citi Bike Demand Analytics - PowerPoint Generator
Creates a professional 6-slide presentation using the existing template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

# Paths
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
EXISTING_PPTX = os.path.join(BASE_DIR, "NYC-Citi-Bike-Demand-A-Big-Data-Exploratory-Analysis-and-Advanced-Interactive-Dashboard.pptx")
OUTPUT_PPTX = os.path.join(BASE_DIR, "NYC-Citi-Bike-Project-Presentation.pptx")
IMAGES_DIR = "/Users/ashishb/.gemini/antigravity/brain/283b26f5-53b6-405b-99e9-04e5e2718a4a"

# Image files
IMAGES = {
    "dashboard": os.path.join(IMAGES_DIR, "dashboard_hero_1765385190938.png"),
    "ml": os.path.join(IMAGES_DIR, "ml_ensemble_1765385210213.png"),
    "nyc": os.path.join(IMAGES_DIR, "nyc_bikeshare_1765385226383.png"),
    "architecture": os.path.join(IMAGES_DIR, "tech_architecture_1765385244325.png"),
    "data": os.path.join(IMAGES_DIR, "data_analysis_1765385261206.png"),
    "future": os.path.join(IMAGES_DIR, "future_vision_1765385275503.png"),
}

# Slide content
SLIDES = [
    {
        "title": "Problem Statement & Dataset",
        "subtitle": "Understanding NYC Urban Mobility Challenges",
        "image": "nyc",
        "bullets": [
            "Challenge: Predict bike demand to optimize station rebalancing",
            "Dataset: 9.9 Million+ trip records (2019-2025)",
            "Coverage: 2,500+ stations across NYC boroughs",
            "Variables: Trip data, weather, temporal patterns, station metadata",
            "Goal: Enable real-time demand forecasting at station level"
        ]
    },
    {
        "title": "Data Processing & Feature Engineering",
        "subtitle": "Big Data Pipeline & 56-Feature Architecture",
        "image": "data",
        "bullets": [
            "Temporal Features: Hour/day cyclical encoding, rush hour flags",
            "Lag Features: 1h, 6h, 24h, 168h historical demand",
            "Rolling Windows: 4h, 12h, 24h, 168h mean/std/min/max",
            "Weather Integration: Temperature, precipitation, wind speed",
            "Advanced: EMA, demand delta, holiday/weekend indicators"
        ]
    },
    {
        "title": "Machine Learning Architecture",
        "subtitle": "Ensemble Model with Gradient Boosting",
        "image": "ml",
        "bullets": [
            "Ensemble: XGBoost + LightGBM + CatBoost weighted voting",
            "Feature Scaling: StandardScaler for tree models",
            "Lazy Loading: Memory-optimized model serving (<1GB RAM)",
            "Prediction: 48-hour rolling forecast per station",
            "Performance: MAE ~2.1, validated on unseen 2025 data"
        ]
    },
    {
        "title": "Interactive Dashboard",
        "subtitle": "Real-Time Analytics & Visualization Platform",
        "image": "dashboard",
        "bullets": [
            "System Overview: KPIs, weather correlation, anomaly detection",
            "Map Explorer: Deck.gl heatmaps, 3D hexagons, 1000+ stations",
            "Prediction Lab: Station-level 48-hour ML forecasts",
            "Station Drilldown: Hourly/daily demand profiles",
            "Premium UI: Glassmorphism, Framer Motion animations"
        ]
    },
    {
        "title": "Technology Stack & Deployment",
        "subtitle": "Cloud-Native Full-Stack Architecture",
        "image": "architecture",
        "bullets": [
            "Backend: FastAPI (Python) - 20+ REST endpoints",
            "Frontend: React + Vite + TailwindCSS + Recharts",
            "Mapping: Deck.gl + MapLibre GL for geospatial viz",
            "Cloud: AWS EC2 with Docker Compose deployment",
            "Data: Parquet files for optimized big data loading"
        ]
    },
    {
        "title": "Results & Future Scope",
        "subtitle": "Impact & Next Steps",
        "image": "future",
        "bullets": [
            "✓ Real-time predictions for 2,500+ NYC bike stations",
            "✓ Weather-aware demand forecasting (temp/precip impact)",
            "✓ Production-ready AWS deployment with 24/7 uptime",
            "→ Future: Reinforcement learning for rebalancing",
            "→ Future: Digital twin simulation, equity analysis"
        ]
    }
]


def set_slide_background_dark(slide, prs):
    """Set dark background for slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 15)  # Very dark blue-black


def add_title_text(slide, title, subtitle, prs):
    """Add title and subtitle text boxes"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.font.size = Pt(18)
    sub_para.font.color.rgb = RGBColor(100, 149, 237)  # Cornflower blue
    sub_para.alignment = PP_ALIGN.LEFT


def add_bullets(slide, bullets, left, top, width, height):
    """Add bullet points"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(220, 220, 230)
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT


def add_image_safe(slide, image_path, left, top, width, height):
    """Add image if it exists"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            return True
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")
    return False


def create_presentation():
    """Create the presentation"""
    print("🚀 Creating NYC Citi Bike Project Presentation...")
    
    # Load existing presentation to get template
    try:
        prs = Presentation(EXISTING_PPTX)
        print(f"✅ Loaded existing template with {len(prs.slides)} slides")
        
        # Keep only slide 1 (title slide), delete rest
        while len(prs.slides) > 1:
            rId = prs.slides._sldIdLst[1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[1]
        print(f"✅ Kept slide 1, removed {len(SLIDES)} slides to replace")
        
    except Exception as e:
        print(f"⚠️ Could not load template: {e}")
        print("Creating new presentation from scratch...")
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
    # Get blank layout (or first layout)
    try:
        blank_layout = prs.slide_layouts[6]  # Usually blank
    except:
        blank_layout = prs.slide_layouts[0]
    
    # Create 6 new slides
    for i, slide_data in enumerate(SLIDES):
        print(f"📄 Creating slide {i+2}: {slide_data['title']}")
        
        slide = prs.slides.add_slide(blank_layout)
        
        # Set dark background
        set_slide_background_dark(slide, prs)
        
        # Add title and subtitle
        add_title_text(slide, slide_data["title"], slide_data["subtitle"], prs)
        
        # Add image on the right side
        image_path = IMAGES.get(slide_data["image"], "")
        if image_path and os.path.exists(image_path):
            add_image_safe(slide, image_path, Inches(6.5), Inches(1.8), Inches(6.5), Inches(5))
        
        # Add bullet points on the left
        add_bullets(slide, slide_data["bullets"], Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
    
    # Save presentation
    prs.save(OUTPUT_PPTX)
    print(f"\n🎉 Presentation saved to: {OUTPUT_PPTX}")
    print(f"   Total slides: {len(prs.slides)}")
    

if __name__ == "__main__":
    create_presentation()
