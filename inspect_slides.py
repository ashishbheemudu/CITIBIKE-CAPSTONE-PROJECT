from pptx import Presentation
import os

SOURCE_PPTX = "/Users/ashishb/Documents/cap/ppt/NYC-Citi-Bike-Demand-A-Big-Data-Exploratory-Analysis-and-Advanced-Interactive-Dashboard.pptx"

if os.path.exists(SOURCE_PPTX):
    prs = Presentation(SOURCE_PPTX)
    print(f"Loaded {SOURCE_PPTX}")
    print(f"Total Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        if i >= 3: break

        print(f"\n--- Slide {i+1} ---")
        print(f"Layout: {slide.slide_layout.name}")
        for shape in slide.shapes:
            print(f"  Shape: {shape.name} (Type: {shape.shape_type})")
            if shape.has_text_frame:
                text = shape.text_frame.text[:50].replace('\n', ' ')
                print(f"    Text: '{text}...'")
