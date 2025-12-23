from pptx import Presentation
import os

SOURCE_PPTX = "/Users/ashishb/Documents/cap/ppt/NYC-Citi-Bike-Project-Presentation.pptx"

if os.path.exists(SOURCE_PPTX):
    prs = Presentation(SOURCE_PPTX)
    print(f"Loaded template from {SOURCE_PPTX}")
    
    print(f"Number of layouts: {len(prs.slide_layouts)}")
    
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout {i}: {layout.name}")
        if len(layout.placeholders) == 0:
            print("  No placeholders.")
        for shape in layout.placeholders:
             print(f"  Placeholder idx {shape.placeholder_format.idx}: type {shape.placeholder_format.type} name '{shape.name}'")
